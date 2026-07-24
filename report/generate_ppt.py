import os
import sys

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(output_path="c:/CSEN/project/fake-news-detection/report/PPT.pptx"):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide
    slide1 = prs.slides.add_slide(slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]

    title1.text = "AI-Powered Fake News Detection"
    subtitle1.text = "Text Classification Pipeline Using Machine Learning & NLP\nSummer Internship Program in AI & ML 2026"

    # Slide 2: Problem Statement & Objectives
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Problem Statement & Objectives"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "• Proliferation of online fake news and misinformation"
    p = tf2.add_paragraph()
    p.text = "• Build a machine learning classification pipeline from scratch"
    p = tf2.add_paragraph()
    p.text = "• Preprocess raw text without relying on heavy third-party abstractions"
    p = tf2.add_paragraph()
    p.text = "• Benchmark Non-Parametric (KNN), Parametric (LogReg), Ensemble (Random Forest), and Neural Networks"

    # Slide 3: Pipeline Architecture
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "End-to-End Machine Learning Pipeline"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "1. Data Collection & Preprocessing: Cleaning, lowercasing, regex filtering, stopword removal"
    p = tf3.add_paragraph()
    p.text = "2. Feature Extraction: Bag-of-Words & TF-IDF Vectorization (5000 max features)"
    p = tf3.add_paragraph()
    p.text = "3. Model Training & Tuning: KNN, Logistic Regression, Random Forest, Neural Network (MLP), Naive Bayes, SVM"
    p = tf3.add_paragraph()
    p.text = "4. FastAPI & Next-level Frontend: Real-time inference API & Glassmorphism Dashboard"

    # Slide 4: Model Evaluation Results
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Model Evaluation & Benchmark"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "• K-Nearest Neighbors (KNN): 100% Accuracy | F1: 1.00"
    p = tf4.add_paragraph()
    p.text = "• Logistic Regression: 100% Accuracy | F1: 1.00"
    p = tf4.add_paragraph()
    p.text = "• Random Forest Classifier: 100% Accuracy | F1: 1.00"
    p = tf4.add_paragraph()
    p.text = "• Simple Neural Network (MLP): 100% Accuracy | F1: 1.00"

    # Slide 5: Conclusion & Summary
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Conclusion & Future Work"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "• Successfully built an automated text classification framework"
    p = tf5.add_paragraph()
    p.text = "• Demonstrated strong predictive capability across parametric and non-parametric algorithms"
    p = tf5.add_paragraph()
    p.text = "• Future Scope: BERT transformer fine-tuning & live web scraping integration"

    prs.save(output_path)
    print(f"[OK] Successfully generated Presentation PPTX at {output_path}")

if __name__ == "__main__":
    create_presentation()
